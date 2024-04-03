import os
import streamlit as st
from langchain.chat_models import ChatOpenAI
from langchain.schema import AIMessage, HumanMessage, SystemMessage
from docx.shared import Pt
from docx import Document
import io
import PyPDF2
from pptx import Presentation
import openpyxl
from PIL import Image
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
 
# Set OpenAI API key
open_api_key = " "#use your openai api key
os.environ["OPENAI_API_KEY"] = open_api_key
 
# Custom prompts for action items
ACTION_ITEM_PROMPTS = {
    "Tabular view": "Extract action items in tabular view from the following meeting transcript:\n",
    "Categorized by priority": "Extract action items categorized by priority (e.g., high, medium, low) from the following meeting transcript:\n",
    "Assigned to specific individuals or roles": "Identify action items assigned to specific individuals or roles from the following meeting transcript:\n"
}
 
# Custom prompt for meeting summarization
MEETING_SUMMARIZATION_PROMPTS = [
    "Generate meeting notes.\n",
    "Summarize meeting Transcript.\n",
    "Key points discussed in the following meeting transcript:\n",
    "Create a table with the ideas discussed and their pros and cons.\n",  
    "Give action items."
]
 
# Default summarization prompt for documents
DEFAULT_DOCUMENT_PROMPT = "Summarize the document and highlight key points:\n"
 
# Add navigation styles
navigation_styles = """
<style>
    /* Navigation bar styles */
    .navbar {
        background-color: #597587;
        padding: 10px;
        margin-bottom: 20px;
    }
 
    .navbar h1 {
        color: #ffffff;
        margin-left: 20px;
        font-size: 24px;
        font-weight: bold;
    }
 
    /* Sidebar styles */
    .sidebar {
        padding: 20px;
        border-radius: 10px;
    }
</style>
"""
 
# Function to process text
def process_text(text, prompt, model_name='gpt-3.5-turbo'):
    messages = [
        SystemMessage(content=f'You are an expert assistant with expertise in summarizing meetings and action items'),
        HumanMessage(content=f'{prompt} {text}')
    ]
    llm_model = ChatOpenAI(model_name=model_name)
    return llm_model(messages).content
 
 
# Function to process Q&A
def process_qa(question, context, model_name='gpt-3.5-turbo', temperature=0.5):
    messages = [
        SystemMessage(content='You are an expert assistant with expertise in answering questions'),
        HumanMessage(content=f'{question} about {context}')
    ]
    llm_model = ChatOpenAI(model_name=model_name)
    return llm_model(messages, temperature=temperature).content
 
 
# Create Streamlit app
st.set_page_config(page_title="Meeting Summarization Demo", page_icon=":memo:")
 
# Apply navigation styles
st.markdown(navigation_styles, unsafe_allow_html=True)
 
# Opening the image
image = Image.open(r'nexturn banner.png')
 
# Displaying the image on the Streamlit app
st.image(image, width=650)
# Main title
st.markdown('<h1 style="font-size: 22px;">👩🏻‍💻 Meeting Summarization</h1>', unsafe_allow_html=True)
 
st.sidebar.title("Select a summarization prompt")
 
# Selectbox for meeting transcript prompt in the sidebar
selected_prompt = st.sidebar.selectbox("Meeting Transcript Prompt", MEETING_SUMMARIZATION_PROMPTS)
 
st.sidebar.markdown(f'<p style="font-size: 18px;">Selected Meeting Prompt: {selected_prompt}</p>', unsafe_allow_html=True)
 
# Sidebar title
st.sidebar.title("Select an action items prompt")
 
# Selectbox for action item prompt in the sidebar
selected_action_prompt = st.sidebar.selectbox("Action Items Prompt", list(ACTION_ITEM_PROMPTS.keys()))
 
# Increase the font size in the sidebar for action items prompt
st.sidebar.markdown(f'<p style="font-size: 18px;">Selected Action Items Prompt: {ACTION_ITEM_PROMPTS[selected_action_prompt]}</p>', unsafe_allow_html=True)
 
st.sidebar.title("Enter Email Address")
 
# Text area for email input in the sidebar
email_input = st.sidebar.text_area("Email:", "")
 
# User input: Meeting transcript or text for action items extraction
meeting_text = st.text_area("Meeting Transcript", "", height=200)
 
# New file uploader for document summarization
document_upload = st.file_uploader("Upload your document for summarization", type=["pdf", "txt", "pptx", "docx", "xlsx"])
 
# Function to send email
def send_email(email, meeting_transcript_result):
    sender_email = "" # Your email address
    receiver_email = email
    password = "" # Your email password
 
    # Combine meeting transcript result and action items result into a single email body
    combined_content = (
        "Meeting Transcript Summarization:\n"
        "Subject: Summary and Action Items from Meeting\n\n"
        "Dear Team,\n\n"
        "I hope this email finds you well. Here's a summary of our recent meeting and the action items discussed:\n\n"+ meeting_transcript_result + "\n\n"
        "Best regards,\n\n"
        "Data Dynamo"
    )
 
    message = MIMEMultipart()
    message["From"] = sender_email
    message["To"] = receiver_email
    message["Subject"] = "Summary and Action Items from Meeting"
 
    message.attach(MIMEText(combined_content, "plain"))
 
    with smtplib.SMTP_SSL("smtp.mail.yahoo.com", 465) as server:
        server.login(sender_email, password)
        server.sendmail(sender_email, receiver_email, message.as_string())
    st.success("Email Sent Successfully!")
 
# Button to trigger meeting transcript summarization
if st.button("Generate Meeting Transcript Result"):
    if meeting_text:
        # Splitting the text into lines and removing empty lines
        lines = meeting_text.split('\n')
 
        # Initialize a variable to store modified text
        raw_text = ''
 
        # Iterate through each line of the transcript
        for line in lines:
            # Check if the line is not empty
            if line.strip():
                # Append the line to the raw_text variable
                raw_text += line.strip() + ' '
        # Process the entire text at once
        result =  process_text(raw_text, selected_prompt)
        st.success("Meeting Transcript Summarization:")
        st.info("Summary: " + result)
 
# Button to trigger action item extraction
if st.button("Extract Action Items"):
    if meeting_text:
        # Splitting the text into lines and removing empty lines
        lines = meeting_text.split('\n')
 
        # Initialize a variable to store modified text
        raw_text = ''
 
        # Iterate through each line of the transcript
        for line in lines:
            # Check if the line is not empty
            if line.strip():
                # Append the line to the raw_text variable
                raw_text += line.strip() + ' '
        action_item_prompt = ACTION_ITEM_PROMPTS[selected_action_prompt]
        action_items_result = process_text(raw_text, action_item_prompt)
 
        st.success("Action Items Extraction:")
        st.info(action_items_result)
 
# Button to trigger document summarization
if st.button("Generate Document Result"):
    if document_upload is not None:
        # Check the file type and extract text accordingly
        if document_upload.type == "application/pdf":
            pdf_reader = PyPDF2.PdfReader(document_upload)
            document_text = ""
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                document_text += page.extract_text()
        elif document_upload.type == "text/plain":
            document_text = document_upload.read().decode("utf-8")
        elif document_upload.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            docx_content = document_upload.read()
            doc = Document(io.BytesIO(docx_content))
            document_text = " ".join([paragraph.text for paragraph in doc.paragraphs])
        elif document_upload.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            pptx_content = document_upload.read()
            presentation = Presentation(io.BytesIO(pptx_content))
            slides_text = []
 
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                slides_text.append(f"Slide {slide_num + 1}:\n{slide_text}")
 
            document_text = "\n".join(slides_text)
        elif document_upload.type == "application/msword":
            doc_content = document_upload.read()
            try:
                doc = Document(io.BytesIO(doc_content))
                paragraphs_text = []
 
                for paragraph_num, paragraph in enumerate(doc.paragraphs):
                    paragraphs_text.append(f"Paragraph {paragraph_num + 1}: {paragraph.text}")
 
                document_text = "\n".join(paragraphs_text)
            except Exception as e:
                st.write(f"Error processing Word document: {e}")
                document_text = None
        elif document_upload.type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            xlsx_content = document_upload.read()
            workbook = openpyxl.load_workbook(io.BytesIO(xlsx_content))
            sheet = workbook.active
            rows_text = []
 
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(map(str, row))
                rows_text.append(row_text)
 
            document_text = "\n".join(rows_text)
        else:
            st.write("Unsupported file type")
            document_text = None
 
        if document_text:
            document_messages = [
                SystemMessage(content='You are an expert assistant with expertise in summarizing text and provide detailed content of the document'),
                HumanMessage(content=f'{DEFAULT_DOCUMENT_PROMPT} for {document_text}')]
            llm_document = ChatOpenAI(model_name='gpt-3.5-turbo')
            document_result = llm_document(document_messages).content
 
            st.success("Document Summarization:")
            st.info(document_result)
 
# Button to trigger meeting transcript summarization and email sending
if st.button("Send Mail"):
    if meeting_text and email_input:
        # Splitting the text into lines and removing empty lines
        lines = meeting_text.split('\n')
 
        # Initialize a variable to store modified text
        raw_text = ''
 
        # Iterate through each line of the transcript
        for line in lines:
            # Check if the line is not empty
            if line.strip():
                # Append the line to the raw_text variable
                raw_text += line.strip() + ' '
 
       
        result = process_text(raw_text, selected_prompt)
 
       
        # Send email
        send_email(email_input, result)
 
# Q&A Section
st.sidebar.title("Q&A Section")
 
# Text area for user to enter a question
qa_question = st.sidebar.text_area("Ask a Question:", "")
 
# Button to trigger Q&A
if st.sidebar.button("Get Answer"):
    if meeting_text and qa_question:
        # Process the question and get the answer
        qa_result = process_qa(qa_question, meeting_text)
        st.sidebar.success("Q&A Result:")
        st.sidebar.info(qa_result)
 
# Close the main content section
st.markdown("</div>", unsafe_allow_html=True)
 